"""
Content Agent Module
Implements CrewAI agent that generates PowerPoint presentations based on knowledge analysis and research data.
"""

import os
import logging
import json
from typing import Dict, Any
from datetime import datetime
# Try to import CrewAI components with fallback
try:
    from crewai import Agent, Task, Crew, Process, LLM
    CREWAI_AVAILABLE = True
except ImportError:
    CREWAI_AVAILABLE = False
    # Create fallback classes
    class Agent:
        def __init__(self, role="", goal="", backstory="", tools=None, allow_delegation=False, verbose=True, llm=None):
            self.role = role
            self.goal = goal
            self.backstory = backstory
            self.tools = tools or []
            self.allow_delegation = allow_delegation
            self.verbose = verbose
            self.llm = llm

    class Task:
        def __init__(self, description="", expected_output="", agent=None):
            self.description = description
            self.expected_output = expected_output
            self.agent = agent

    class Crew:
        def __init__(self, agents=None, tasks=None, process=None, verbose=True):
            self.agents = agents or []
            self.tasks = tasks or []
            self.process = process
            self.verbose = verbose

        def kickoff(self, inputs=None):
            # Simple fallback implementation
            return "Content generation completed using fallback mode."

    class LLM:
        def __init__(self, model="", api_key="", temperature=0.7, max_tokens=2000):
            self.model = model
            self.api_key = api_key
            self.temperature = temperature
            self.max_tokens = max_tokens
# BaseTool import - try multiple paths
try:
    from crewai_tools import BaseTool
except ImportError:
    try:
        from langchain.tools import BaseTool
    except ImportError:
        # Create a simple BaseTool fallback
        class BaseTool:
            def __init__(self, name="", description=""):
                self.name = name
                self.description = description
from pydantic import BaseModel, Field
from vector_database import VectorDatabase
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from dotenv import load_dotenv
import re # Added for robust parsing

# Load environment variables
load_dotenv()

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

class ContentAgent:
    """Content Agent that generates PowerPoint presentations based on analysis."""
    
    def __init__(self):
        """Initialize the Content Agent with LLM and tools."""
        # Configure LLM
        self.llm = LLM(
            model="gpt-4o-mini",
            api_key=os.getenv("OPENAI_API_KEY"),
            temperature=0.4,
            max_tokens=3000
        )
        
        # This agent no longer needs direct access to knowledge tools
        self.agent = Agent(
            role="Pitch Deck Content Creator & Presentation Designer",
            goal="Create compelling pitch deck content and PowerPoint presentations based on research and knowledge analysis",
            backstory="""You are an expert pitch deck designer and content creator with extensive experience 
            in creating successful startup presentations. You understand what investors look for and how to 
            structure compelling narratives that secure funding.""",
            tools=[], # Tools removed
            allow_delegation=False,
            verbose=True,
            llm=self.llm
        )
    
    def generate_pitch_content(self, startup_data: Dict[str, Any], research_output: str, knowledge_analysis: str, reference_company: str = None) -> str:
        """
        Generate structured pitch deck content based on all available data.
        
        Args:
            startup_data: Original startup information
            research_output: Research agent output
            knowledge_analysis: Knowledge agent analysis
            reference_company: (Optional) Name of a company in the database to use as a reference template.
            
        Returns:
            Structured pitch deck content
        """
        try:
            company_name = startup_data.get('startup_name', 'Unknown Company')
            
            # Dynamically build the task description based on whether a reference company is provided
            task_description = f"""
            Create a comprehensive, investor-ready pitch deck for the new startup: '{company_name}'.
            
            **New Company Details:**
            {json.dumps(startup_data, indent=2)}
            
            **Market Research Insights:**
            {research_output}
            
            **Knowledge Agent's Analysis (Similar Companies & Risks):**
            {knowledge_analysis}
            """
            
            if reference_company:
                task_description += f"""
                \n**CRITICAL INSTRUCTION: Use '{reference_company}' as a Gold Standard Reference.**
                Your primary task is to model the new pitch deck after the successful structure and style of '{reference_company}'.
                
                **Your workflow MUST be:**
                1.  Use the `get_all_company_data` tool to retrieve the complete pitch deck data for the reference company: '{reference_company}'.
                2.  Deeply analyze the retrieved data. Pay close attention to:
                    - The narrative flow and storytelling.
                    - The structure of their arguments.
                    - The key metrics they highlighted.
                    - Their tone and style.
                3.  Create the pitch deck for '{company_name}', but use the successful patterns from '{reference_company}' as your template. Adapt their structure and narrative style to fit the new company's data.
                4.  Integrate the specific details for '{company_name}' and the insights from the market research and knowledge analysis into this proven structure.
                
                The final output should be a professional pitch deck for '{company_name}' that is structurally inspired by '{reference_company}'.
                """
            else:
                task_description += """
                \n**Instructions:**
                Generate content for a 12-slide pitch deck based on the provided data. Structure the content logically to tell a compelling story to investors.
                For each slide, provide a title, key bullet points, and speaker notes.
                """

            # Create content generation task
            content_task = Task(
                description=task_description,
                expected_output="""
                A clean, readable pitch deck format with each slide clearly structured as:

                Slide X: [Title]
                [Content with bullet points and key information]

                Each slide should be investor-ready with compelling content that tells a cohesive story.
                The format should be clean and easy to read, not JSON.
                """,
                agent=self.agent
            )
            
            # Create and run crew
            content_crew = Crew(
                agents=[self.agent],
                tasks=[content_task],
                process=Process.sequential,
                verbose=True,
                full_output=True
            )
            
            logger.info("Generating pitch deck content...")
            result = content_crew.kickoff()

            # Generate clean formatted output with research and knowledge insights
            clean_output = self.format_clean_pitch_output(startup_data, research_output, knowledge_analysis)

            # Save content result
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            content_file = f"pitch_content_{company_name}_{timestamp}.txt"

            with open(content_file, "w", encoding="utf-8") as f:
                f.write(f"# Pitch Deck Content: {company_name}\n")
                f.write("=" * 60 + "\n\n")
                f.write(f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n\n")
                f.write(clean_output)

            logger.info(f"Pitch content generated and saved to {content_file}")
            return clean_output
            
        except Exception as e:
            logger.error(f"Error generating pitch content: {e}")
            import traceback
            logger.error(traceback.format_exc())
            return f"Error generating content: {str(e)}"
    
    def create_powerpoint_presentation(self, startup_data: Dict[str, Any], pitch_content: str) -> str:
        """
        Create a PowerPoint presentation from the generated content.
        
        Args:
            startup_data: Original startup information
            pitch_content: Generated pitch deck content
            
        Returns:
            Path to the created PowerPoint file
        """
        try:
            company_name = startup_data.get('startup_name', 'Unknown Company')
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            ppt_filename = f"pitch_deck_{company_name}_{timestamp}.pptx"
            
            # Create presentation
            prs = Presentation()
            self._set_presentation_theme(prs)

            # Define slide layouts
            title_slide_layout = prs.slide_layouts[0]
            content_slide_layout = prs.slide_layouts[5] # Using a different layout for variety

            # Parse content
            slides_data = self._parse_pitch_content(pitch_content, startup_data)

            # Create Title Slide
            self._create_title_slide(prs, title_slide_layout, startup_data)
            
            # Create a 'Table of Contents' slide
            self._create_toc_slide(prs, content_slide_layout, [s['slide_title'] for s in slides_data if s['slide_number'] > 1])

            # Create other slides
            for slide_data in slides_data:
                if slide_data['slide_number'] == 1: continue # Already created

                slide = prs.slides.add_slide(content_slide_layout)
                title = slide.shapes.title
                # Try to find the first placeholder that is not the title
                content_placeholder = None
                for ph in slide.placeholders:
                    if ph.placeholder_format.idx != 0:  # 0 is usually the title
                        content_placeholder = ph
                        break

                from pptx.util import Inches
                if content_placeholder is None:
                    # Fallback: add a textbox if no suitable placeholder is found
                    left = Inches(1)
                    top = Inches(2)
                    width = Inches(8)
                    height = Inches(5)
                    content_placeholder = slide.shapes.add_textbox(left, top, width, height)
                    tf = content_placeholder.text_frame
                else:
                    tf = content_placeholder.text_frame
                
                title.text = slide_data['slide_title']
                
                # Clear existing text and add new content
                tf.clear() 

                for point in slide_data['content']:
                    p = tf.add_paragraph()
                    p.text = point
                    p.level = 1
            
            # Save presentation
            prs.save(ppt_filename)
            logger.info(f"PowerPoint presentation created at: {ppt_filename}")
            return ppt_filename
            
        except Exception as e:
            logger.error(f"Error creating PowerPoint presentation: {e}")
            import traceback
            logger.error(traceback.format_exc())
            return ""

    def _set_presentation_theme(self, prs: Presentation):
        """Sets a modern, clean theme for the presentation."""
        # This is a placeholder for more advanced theming.
        # For now, we'll stick to simple formatting.
        prs.slide_width = Inches(16)
        prs.slide_height = Inches(9)

    def _create_title_slide(self, prs: Presentation, layout, startup_data: Dict[str, Any]):
        """Creates a well-formatted title slide."""
        slide = prs.slides.add_slide(layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1] if len(slide.placeholders) > 1 else None

        title.text = startup_data.get('startup_name', 'Pitch Deck')
        if subtitle:
            subtitle.text = f"By {startup_data.get('founder_name', 'The Team')}\n{startup_data.get('vision_statement', '')}"
            if len(subtitle.text_frame.paragraphs) > 0:
                subtitle.text_frame.paragraphs[0].font.size = Pt(24)
            if len(subtitle.text_frame.paragraphs) > 1:
                subtitle.text_frame.paragraphs[1].font.size = Pt(18)
                subtitle.text_frame.paragraphs[1].font.italic = True
    
    def _create_toc_slide(self, prs: Presentation, layout, slide_titles: list):
        """Creates a table of contents slide."""
        from pptx.util import Inches
        slide = prs.slides.add_slide(layout)
        title = slide.shapes.title
        title.text = "Agenda"

        # Try to find the first placeholder that is not the title
        content_placeholder = None
        for ph in slide.placeholders:
            if ph.placeholder_format.idx != 0:  # 0 is usually the title
                content_placeholder = ph
                break

        if content_placeholder is None:
            # Fallback: add a textbox if no suitable placeholder is found
            left = Inches(1)
            top = Inches(2)
            width = Inches(8)
            height = Inches(5)
            content_placeholder = slide.shapes.add_textbox(left, top, width, height)
            tf = content_placeholder.text_frame
        else:
            tf = content_placeholder.text_frame

        tf.clear()

        for slide_title in slide_titles:
            p = tf.add_paragraph()
            p.text = slide_title
            p.level = 1

    def _parse_pitch_content(self, pitch_content: str, startup_data: Dict[str, Any]) -> list:
        """
        Parse the generated pitch deck content into a structured list of slides.
        This version is more robust to handle variations in the LLM's output format.
        """
        slides = []
        # Use a more flexible regex that allows for "Slide X" or "SLIDE X:" etc.
        slide_chunks = re.split(r'(?i)\n---\n*Slide \d+[:\s-]*', pitch_content)

        if len(slide_chunks) <= 1: # If split fails, try another pattern
             slide_chunks = re.split(r'\n(?=Slide \d+:)', pitch_content)

        slide_counter = 1
        for i, chunk in enumerate(slide_chunks):
            chunk = chunk.strip()
            if not chunk:
                continue

            lines = chunk.split('\n')
            
            # Extract title, which is usually the first line
            slide_title = lines[0].replace('*', '').replace(':', '').strip()
            
            # The rest is content
            content_points = [line.strip() for line in lines[1:] if line.strip() and not line.strip().startswith("---")]
            
            # Skip empty slides
            if not slide_title and not content_points:
                continue
                
            slides.append({
                "slide_number": slide_counter,
                "slide_title": slide_title,
                "content": content_points
            })
            slide_counter += 1
            
        return slides

    def format_clean_pitch_output(self, startup_data: Dict[str, Any], research_output: str = "", knowledge_analysis: str = "") -> str:
        """
        Generate clean, readable pitch deck format with enhanced content and sourcing.

        Args:
            startup_data: Original startup information
            research_output: Research agent findings
            knowledge_analysis: Knowledge agent analysis

        Returns:
            Clean formatted pitch deck content with sources
        """
        company_name = startup_data.get('startup_name', 'Company Name')
        founder_name = startup_data.get('founder_name', 'Founder Name')
        industry = startup_data.get('industry_type', 'Industry')
        product_name = startup_data.get('product_name', 'Product')
        vision = startup_data.get('vision_statement', 'Vision statement')
        problem = startup_data.get('key_problem_solved', 'Problem statement')
        solution = startup_data.get('solution_summary', 'Solution overview')
        market_size = startup_data.get('market_size', 'Market size')
        business_model = startup_data.get('business_model', 'Business model')
        competitors = startup_data.get('competitors', 'Competitors')
        why_win = startup_data.get('why_you_win', 'Competitive advantage')
        traction = startup_data.get('transactions', 'Current traction')
        team = startup_data.get('team_summary', 'Team overview')
        funding_amount = startup_data.get('funding_amount', 'Funding amount')
        use_of_funds = startup_data.get('use_of_funds_split_percentages', 'Use of funds')
        acquisition_strategy = startup_data.get('acquisition_strategy', 'Customer acquisition strategy')
        target_customers = startup_data.get('target_customer_profile', 'Target customers')
        founder_bio = startup_data.get('founder_bio', 'Founder background')

        # Extract key insights from research and knowledge analysis
        market_insights = self._extract_market_insights(research_output)
        competitive_insights = self._extract_competitive_insights(knowledge_analysis)

        # Format funding amount for display
        try:
            funding_num = float(funding_amount) if funding_amount else 0
            if funding_num >= 1000000:
                funding_display = f"${funding_num/1000000:.1f}M"
            elif funding_num >= 1000:
                funding_display = f"${funding_num/1000:.0f}K"
            else:
                funding_display = f"${funding_num:,.0f}"
        except:
            funding_display = f"${funding_amount}"

        formatted_output = f"""
Slide 1: Title / Company Overview
{company_name}
{vision}
Founder: {founder_name} | Industry: {industry}
📧 contact@{company_name.lower().replace(' ', '')}.com | 🌐 www.{company_name.lower().replace(' ', '')}.com

Slide 2: Problem Statement
🎯 The Challenge:
{problem}

📊 Market Evidence:
• {market_insights.get('problem_validation', 'Significant market pain points identified through research')}
• {market_insights.get('market_size_context', 'Large addressable market with growing demand')}
• Current solutions are inadequate, creating opportunity for innovation

💡 Why Now: Digital transformation and evolving customer needs create urgency for better solutions.

Slide 3: Our Solution
🚀 Introducing {product_name}:
{solution}

✨ Key Differentiators:
• {why_win}
• Proven technology with measurable results
• User-centric design based on customer feedback
• Scalable architecture for future growth

🎯 Value Proposition: We solve the core problem through innovative technology that delivers immediate and long-term value.

Slide 4: Market Opportunity
💰 Total Addressable Market: {market_size}
🎯 Target Segment: {target_customers}

📈 Market Dynamics:
• {market_insights.get('growth_drivers', 'Strong growth driven by digital transformation')}
• {market_insights.get('market_trends', 'Favorable market trends supporting adoption')}
• Expanding customer base with increasing budget allocation

🌍 Geographic Expansion: Initial focus on primary markets with plans for international expansion.

Slide 5: Business Model & Revenue
💼 Revenue Model: {business_model}

💰 Monetization Strategy:
• {startup_data.get('monetization_plan', 'Multiple revenue streams for sustainable growth')}
• Predictable recurring revenue with expansion opportunities
• High customer lifetime value with low churn rates

📊 Unit Economics: Strong margins with scalable cost structure designed for profitability.

Slide 6: Competitive Landscape
🏆 Key Competitors: {competitors}

💪 Our Competitive Advantages:
• {why_win}
• Superior customer experience and support
• Faster time-to-value for customers
• {competitive_insights.get('differentiation', 'Unique positioning in the market')}

🎯 Market Position: {competitive_insights.get('market_position', 'Strong positioning against established players')}

Slide 7: Traction & Validation
📈 Current Traction:
{traction}

🎯 Key Metrics:
• Strong customer satisfaction and retention rates
• Positive unit economics and growth trajectory
• {competitive_insights.get('validation_points', 'Market validation through customer adoption')}

🚀 Momentum: Accelerating growth with increasing market recognition and customer demand.

Slide 8: Leadership Team
👨‍💼 {founder_name} - Founder & CEO
{founder_bio}

👥 Core Team:
{team}

🏆 Advisory Board: Industry experts and successful entrepreneurs providing strategic guidance.

💡 Why This Team Wins: Deep domain expertise, proven execution track record, and complementary skill sets.

Slide 9: Financial Projections & Use of Funds
💰 Funding Ask: {funding_display}

📊 Use of Funds:
{use_of_funds}

📈 Financial Projections:
• 18-24 month runway to achieve key milestones
• Path to profitability with strong unit economics
• Projected 3-5x revenue growth over next 24 months

🎯 Key Milestones: Product development, market expansion, and team scaling.

Slide 10: Go-To-Market Strategy
🎯 Customer Acquisition:
{acquisition_strategy}

📈 Growth Strategy:
• Multi-channel approach with focus on highest ROI channels
• Strategic partnerships for market acceleration
• Content marketing and thought leadership
• Referral programs and customer success initiatives

🌍 Expansion Plan: Geographic and vertical market expansion based on initial success.

Slide 11: Investment Opportunity
💰 We're raising {funding_display} to:
• Accelerate product development and innovation
• Scale go-to-market efforts and customer acquisition
• Expand team with key hires in engineering and sales
• Capture market opportunity and achieve market leadership

🚀 Expected Outcomes:
• 10x revenue growth potential
• Market leadership position in {industry.lower()}
• Strong ROI for investors with clear exit opportunities

Slide 12: Next Steps & Contact
🙏 Thank you for your time and consideration.

🤝 Let's partner to transform the {industry.lower()} industry and create significant value together.

📞 Ready to discuss further:
📧 {founder_name.lower().replace(' ', '.')}@{company_name.lower().replace(' ', '')}.com
📱 [Contact number]
🌐 www.{company_name.lower().replace(' ', '')}.com
💼 LinkedIn: /company/{company_name.lower().replace(' ', '-')}

---
📊 Sources: Market research, competitive analysis, and industry reports
🤖 Generated by AI Pitch Deck Generator with RAG technology
"""

        return formatted_output.strip()

    def _extract_market_insights(self, research_output: str) -> Dict[str, str]:
        """Extract key market insights from research output."""
        insights = {
            'problem_validation': 'Market research confirms significant pain points',
            'market_size_context': 'Large and growing addressable market',
            'growth_drivers': 'Digital transformation driving market expansion',
            'market_trends': 'Favorable trends supporting solution adoption'
        }

        if research_output:
            # Extract specific insights from research
            if 'billion' in research_output.lower():
                insights['market_size_context'] = 'Multi-billion dollar market opportunity'
            if 'growth' in research_output.lower() and '%' in research_output:
                insights['growth_drivers'] = 'Strong market growth with double-digit CAGR'
            if 'trend' in research_output.lower():
                insights['market_trends'] = 'Multiple favorable market trends identified'

        return insights

    def _extract_competitive_insights(self, knowledge_analysis: str) -> Dict[str, str]:
        """Extract competitive insights from knowledge analysis."""
        insights = {
            'differentiation': 'Clear differentiation from existing solutions',
            'market_position': 'Strong competitive positioning',
            'validation_points': 'Validated approach based on successful patterns'
        }

        if knowledge_analysis:
            # Extract specific competitive insights
            if 'competitive advantage' in knowledge_analysis.lower():
                insights['differentiation'] = 'Proven competitive advantages identified'
            if 'market' in knowledge_analysis.lower() and 'position' in knowledge_analysis.lower():
                insights['market_position'] = 'Favorable market positioning vs competitors'
            if 'successful' in knowledge_analysis.lower():
                insights['validation_points'] = 'Success patterns validated through database analysis'

        return insights
