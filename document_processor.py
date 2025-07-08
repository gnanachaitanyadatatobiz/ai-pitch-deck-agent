"""
Document Processing Module for Pitch Deck Analysis
Extracts text from PDFs and PowerPoint files, processes them into chunks for vector storage.
"""

import os
import logging
import re
from typing import List, Dict, Any, Tuple
from pathlib import Path
import PyPDF2
import pdfplumber
from pptx import Presentation
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.schema import Document
import hashlib

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

class DocumentProcessor:
    """Handles extraction and processing of documents from the Data folder."""
    
    def __init__(self, data_folder: str = "Data", chunk_size: int = 1000, chunk_overlap: int = 200):
        """
        Initialize the document processor.
        
        Args:
            data_folder: Path to the folder containing documents
            chunk_size: Size of text chunks for vector storage
            chunk_overlap: Overlap between chunks
        """
        self.data_folder = Path(data_folder)
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=chunk_size,
            chunk_overlap=chunk_overlap,
            length_function=len,
            separators=["\n\n", "\n", " ", ""]
        )
        
    def extract_text_from_pdf(self, file_path: Path) -> Tuple[str, Dict[str, Any]]:
        """
        Extract text from PDF file using multiple methods for better accuracy.
        
        Args:
            file_path: Path to the PDF file
            
        Returns:
            Tuple of (extracted_text, metadata)
        """
        text = ""
        metadata = {
            "file_name": file_path.name,
            "file_type": "pdf",
            "file_path": str(file_path),
            "pages": 0
        }
        
        try:
            # Try pdfplumber first (better for complex layouts)
            with pdfplumber.open(file_path) as pdf:
                metadata["pages"] = len(pdf.pages)
                for page_num, page in enumerate(pdf.pages, 1):
                    page_text = page.extract_text()
                    if page_text:
                        text += f"\n--- Page {page_num} ---\n{page_text}\n"
                        
        except Exception as e:
            logger.warning(f"pdfplumber failed for {file_path}, trying PyPDF2: {e}")
            
            # Fallback to PyPDF2
            try:
                with open(file_path, 'rb') as file:
                    pdf_reader = PyPDF2.PdfReader(file)
                    metadata["pages"] = len(pdf_reader.pages)
                    
                    for page_num, page in enumerate(pdf_reader.pages, 1):
                        page_text = page.extract_text()
                        if page_text:
                            text += f"\n--- Page {page_num} ---\n{page_text}\n"
                            
            except Exception as e2:
                logger.error(f"Both PDF extraction methods failed for {file_path}: {e2}")
                return "", metadata
        
        # Clean and normalize text
        text = self._clean_text(text)
        metadata["character_count"] = len(text)
        metadata["word_count"] = len(text.split())
        
        return text, metadata
    
    def extract_text_from_pptx(self, file_path: Path) -> Tuple[str, Dict[str, Any]]:
        """
        Extract text from PowerPoint file.
        
        Args:
            file_path: Path to the PPTX file
            
        Returns:
            Tuple of (extracted_text, metadata)
        """
        text = ""
        metadata = {
            "file_name": file_path.name,
            "file_type": "pptx",
            "file_path": str(file_path),
            "slides": 0
        }
        
        try:
            presentation = Presentation(file_path)
            metadata["slides"] = len(presentation.slides)
            
            for slide_num, slide in enumerate(presentation.slides, 1):
                slide_text = f"\n--- Slide {slide_num} ---\n"
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_text += shape.text + "\n"
                    
                    # Extract text from tables
                    if shape.has_table:
                        for row in shape.table.rows:
                            row_text = []
                            for cell in row.cells:
                                if cell.text:
                                    row_text.append(cell.text.strip())
                            if row_text:
                                slide_text += " | ".join(row_text) + "\n"
                
                text += slide_text
                
        except Exception as e:
            logger.error(f"Failed to extract text from {file_path}: {e}")
            return "", metadata
        
        # Clean and normalize text
        text = self._clean_text(text)
        metadata["character_count"] = len(text)
        metadata["word_count"] = len(text.split())
        
        return text, metadata
    
    def _clean_text(self, text: str) -> str:
        """
        Clean and normalize extracted text.
        
        Args:
            text: Raw extracted text
            
        Returns:
            Cleaned text
        """
        if not text:
            return ""
        
        # Remove excessive whitespace
        text = re.sub(r'\s+', ' ', text)
        
        # Remove special characters but keep basic punctuation
        text = re.sub(r'[^\w\s\.\,\!\?\;\:\-\(\)\[\]\"\'\/\%\$\&\@\#]', ' ', text)
        
        # Remove excessive newlines but preserve paragraph breaks
        text = re.sub(r'\n\s*\n', '\n\n', text)
        
        # Strip leading/trailing whitespace
        text = text.strip()
        
        return text
    
    def extract_company_info(self, text: str, file_name: str) -> Dict[str, str]:
        """
        Extract company information from the document text.
        
        Args:
            text: Document text
            file_name: Name of the file
            
        Returns:
            Dictionary with extracted company information
        """
        company_info = {
            "company_name": "",
            "industry": "",
            "stage": "",
            "funding_amount": "",
            "year": ""
        }
        
        # Extract company name from filename (common pattern)
        name_patterns = [
            r'^([a-zA-Z]+)',  # First word in filename
            r'([A-Z][a-z]+)',  # Capitalized words
        ]
        
        for pattern in name_patterns:
            match = re.search(pattern, file_name)
            if match:
                company_info["company_name"] = match.group(1).lower().capitalize()
                break
        
        # Extract year from filename
        year_match = re.search(r'(20\d{2})', file_name)
        if year_match:
            company_info["year"] = year_match.group(1)
        
        # Extract funding information from text
        funding_patterns = [
            r'\$(\d+(?:\.\d+)?)\s*(?:million|M|mil)',
            r'\$(\d+(?:\.\d+)?)\s*(?:billion|B|bil)',
            r'(\d+(?:\.\d+)?)\s*(?:million|M|mil)',
            r'raise\s*\$?(\d+(?:\.\d+)?)',
            r'funding\s*\$?(\d+(?:\.\d+)?)'
        ]
        
        for pattern in funding_patterns:
            match = re.search(pattern, text, re.IGNORECASE)
            if match:
                company_info["funding_amount"] = match.group(1)
                break
        
        return company_info

    def process_document(self, file_path: Path) -> List[Document]:
        """
        Process a single document and return chunks as LangChain Documents.

        Args:
            file_path: Path to the document

        Returns:
            List of Document objects with text chunks and metadata
        """
        documents = []

        try:
            # Determine file type and extract text
            if file_path.suffix.lower() == '.pdf':
                text, metadata = self.extract_text_from_pdf(file_path)
            elif file_path.suffix.lower() in ['.pptx', '.ppt']:
                text, metadata = self.extract_text_from_pptx(file_path)
            else:
                logger.warning(f"Unsupported file type: {file_path}")
                return documents

            if not text.strip():
                logger.warning(f"No text extracted from {file_path}")
                return documents

            # Extract company information
            company_info = self.extract_company_info(text, file_path.name)
            metadata.update(company_info)

            # Generate unique document ID
            doc_id = hashlib.md5(str(file_path).encode()).hexdigest()
            metadata["document_id"] = doc_id

            # Split text into chunks
            text_chunks = self.text_splitter.split_text(text)

            # Create Document objects for each chunk
            for i, chunk in enumerate(text_chunks):
                chunk_metadata = metadata.copy()
                chunk_metadata["chunk_id"] = f"{doc_id}_{i}"
                chunk_metadata["chunk_index"] = i
                chunk_metadata["total_chunks"] = len(text_chunks)

                doc = Document(
                    page_content=chunk,
                    metadata=chunk_metadata
                )
                documents.append(doc)

            logger.info(f"Processed {file_path.name}: {len(text_chunks)} chunks created")

        except Exception as e:
            logger.error(f"Error processing {file_path}: {e}")

        return documents

    def process_all_documents(self) -> List[Document]:
        """
        Process all documents in the data folder.

        Returns:
            List of all Document objects from all processed files
        """
        all_documents = []

        if not self.data_folder.exists():
            logger.error(f"Data folder {self.data_folder} does not exist")
            return all_documents

        # Get all PDF and PowerPoint files
        file_patterns = ['*.pdf', '*.pptx', '*.ppt']
        files_to_process = []

        for pattern in file_patterns:
            files_to_process.extend(self.data_folder.glob(pattern))

        logger.info(f"Found {len(files_to_process)} files to process")

        # Process each file
        for file_path in files_to_process:
            logger.info(f"Processing: {file_path.name}")
            documents = self.process_document(file_path)
            all_documents.extend(documents)

        logger.info(f"Total documents created: {len(all_documents)}")
        return all_documents

    def get_processing_summary(self, documents: List[Document]) -> Dict[str, Any]:
        """
        Generate a summary of the processing results.

        Args:
            documents: List of processed documents

        Returns:
            Summary statistics
        """
        if not documents:
            return {"total_documents": 0}

        # Group by file
        files = {}
        companies = set()

        for doc in documents:
            file_name = doc.metadata.get("file_name", "unknown")
            company_name = doc.metadata.get("company_name", "")

            if file_name not in files:
                files[file_name] = {
                    "chunks": 0,
                    "total_chars": 0,
                    "company": company_name,
                    "file_type": doc.metadata.get("file_type", "unknown")
                }

            files[file_name]["chunks"] += 1
            files[file_name]["total_chars"] += len(doc.page_content)

            if company_name:
                companies.add(company_name.lower())

        summary = {
            "total_documents": len(documents),
            "total_files": len(files),
            "unique_companies": len(companies),
            "companies": sorted(list(companies)),
            "files": files,
            "average_chunk_size": sum(len(doc.page_content) for doc in documents) / len(documents)
        }

        return summary
